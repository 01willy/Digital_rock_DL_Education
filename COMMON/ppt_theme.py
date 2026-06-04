"""
공통 PPT 테마 빌더 — 사용자 디자인(Pretendard + 오렌지)을 차용한 헬퍼.

사용 예:
    from ppt_theme import new_presentation, add_title_slide, add_content_slide, \\
        ORANGE, BODY, SECONDARY, FONT, FONT_BOLD, FONT_BLACK
    prs = new_presentation()
    add_title_slide(prs, "Week 1", "Digital Rock 데이터 이해", "1조 (Advanced)")
    s = add_content_slide(prs, "1. Digital Rock이란?", page_num=2, total=18)
    ...
    prs.save("out.pptx")
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ──────────────────────────────────────────────────────────────
# 색상 (사용자 디자인 차용 — #EA851B 오렌지 메인)
# ──────────────────────────────────────────────────────────────
ORANGE = RGBColor(0xEA, 0x85, 0x1B)
ORANGE_ALT = RGBColor(0xF1, 0x80, 0x1F)
BODY = RGBColor(0x18, 0x18, 0x18)
SECONDARY = RGBColor(0x7A, 0x7A, 0x7A)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
RED = RGBColor(0xC8, 0x3A, 0x3A)
CREAM = RGBColor(0xFA, 0xF5, 0xEE)
LIGHT_GRAY = RGBColor(0xF6, 0xF6, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = 'Pretendard'
FONT_MED = 'Pretendard Medium'
FONT_SEMI = 'Pretendard SemiBold'
FONT_BOLD = 'Pretendard ExtraBold'
FONT_BLACK = 'Pretendard Black'


def new_presentation():
    """16:9 widescreen presentation 생성."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _set_run(run, text, size, font=FONT, color=BODY, bold=False):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, text, left, top, width, height,
             size=20, font=FONT_SEMI, color=BODY, bold=False,
             align=PP_ALIGN.LEFT, line_spacing=1.15):
    """텍스트 박스 추가. text가 list면 각 항목이 한 paragraph."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    lines = text if isinstance(text, list) else text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        _set_run(run, line, size, font=font, color=color, bold=bold)
    return tb


def add_bullets(slide, bullets, left, top, width, height,
                size=18, font=FONT_SEMI, color=BODY, bullet_char='•',
                line_spacing=1.3):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        # bullet char in orange
        run1 = p.add_run()
        _set_run(run1, f'{bullet_char}  ', size, font=FONT_BOLD, color=ORANGE)
        # body in body color
        run2 = p.add_run()
        _set_run(run2, b, size, font=font, color=color)
    return tb


def add_footer(slide, page_num, total, label='Week 1'):
    """좌하단 페이지 정보 + 우하단 페이지 번호."""
    add_text(slide, label, Inches(0.4), Inches(7.1), Inches(8), Inches(0.3),
             size=10, font=FONT_MED, color=SECONDARY)
    add_text(slide, f'{page_num} / {total}', Inches(12.2), Inches(7.1),
             Inches(1), Inches(0.3),
             size=10, font=FONT_MED, color=SECONDARY, align=PP_ALIGN.RIGHT)


def add_title_slide(prs, eyebrow, title, subtitle=None):
    """표지 슬라이드 — 흰 배경, 좌측 오렌지 띠, 큰 타이틀."""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)

    # 좌측 세로 오렌지 띠
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), prs.slide_height)
    bar.line.fill.background()
    bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE
    bar.text_frame.text = ''

    # Eyebrow (작은 안내)
    add_text(s, eyebrow, Inches(0.8), Inches(2.0), Inches(11), Inches(0.5),
             size=20, font=FONT_BOLD, color=ORANGE)
    # Main title
    add_text(s, title, Inches(0.8), Inches(2.6), Inches(11), Inches(1.5),
             size=48, font=FONT_BLACK, color=BODY)
    # Subtitle
    if subtitle:
        add_text(s, subtitle, Inches(0.8), Inches(4.3), Inches(11), Inches(0.6),
                 size=22, font=FONT_SEMI, color=SECONDARY)
    # 하단 작은 라벨
    add_text(s, 'Digital Rock — Sparse Slice Interpolation · 6주 학생 코스',
             Inches(0.8), Inches(6.8), Inches(11), Inches(0.4),
             size=12, font=FONT_MED, color=SECONDARY)
    return s


def add_content_slide(prs, title, page_num=None, total=None, label='Week 1'):
    """본문 슬라이드 베이스 — 상단 타이틀 + 오렌지 underline + 푸터."""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)

    # Title
    add_text(s, title, Inches(0.5), Inches(0.35), Inches(12.5), Inches(0.7),
             size=28, font=FONT_BLACK, color=ORANGE)
    # Orange underline (가는 띠)
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.55), Inches(1.05),
                              Inches(1.2), Inches(0.06))
    line.line.fill.background()
    line.fill.solid(); line.fill.fore_color.rgb = ORANGE

    if page_num is not None and total is not None:
        add_footer(s, page_num, total, label=label)
    return s


def add_section_slide(prs, section_no, section_title, page_num=None, total=None, label='Week 1'):
    """챕터 구분 슬라이드 — 크림 배경 + 큰 챕터 번호."""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)

    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = CREAM
    bg.text_frame.text = ''

    add_text(s, f'Section {section_no}', Inches(0.8), Inches(2.5),
             Inches(11), Inches(0.6),
             size=20, font=FONT_BOLD, color=ORANGE)
    add_text(s, section_title, Inches(0.8), Inches(3.1),
             Inches(11), Inches(1.5),
             size=44, font=FONT_BLACK, color=BODY)
    if page_num is not None and total is not None:
        add_footer(s, page_num, total, label=label)
    return s


def add_box(slide, left, top, width, height, fill=CREAM, border=None):
    """배경 박스 (rounded rectangle)."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = fill
    if border is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = border
        box.line.width = Pt(1.5)
    box.text_frame.text = ''
    return box


def add_picture(slide, path, left, top, width=None, height=None):
    """이미지 추가."""
    kwargs = {}
    if width is not None: kwargs['width'] = width
    if height is not None: kwargs['height'] = height
    return slide.shapes.add_picture(path, left, top, **kwargs)


def add_code_block(slide, code, left, top, width, height, size=14):
    """코드 블록 — 회색 배경 + monospace."""
    add_box(slide, left, top, width, height, fill=LIGHT_GRAY)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.1)
    tf.margin_right = Inches(0.1); tf.margin_bottom = Inches(0.1)
    lines = code.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.2
        run = p.add_run()
        _set_run(run, line, size, font='Consolas', color=BODY)
    return tb


def add_qa_pair(slide, q, a, top, q_size=18, a_size=16):
    """질문 → 답 패턴 (한 줄)."""
    add_text(slide, q, Inches(0.7), top, Inches(12), Inches(0.5),
             size=q_size, font=FONT_BOLD, color=ORANGE)
    add_text(slide, a, Inches(1.0), top + Inches(0.55), Inches(12), Inches(0.8),
             size=a_size, font=FONT_SEMI, color=BODY)


def add_two_column(slide, left_title, left_bullets, right_title, right_bullets,
                   top=Inches(1.4), size=18):
    """두 칼럼 비교 슬라이드."""
    col_w = Inches(6.0)
    add_text(slide, left_title, Inches(0.6), top, col_w, Inches(0.5),
             size=22, font=FONT_BLACK, color=ORANGE)
    add_bullets(slide, left_bullets, Inches(0.6), top + Inches(0.6),
                col_w, Inches(5), size=size)
    add_text(slide, right_title, Inches(7.0), top, col_w, Inches(0.5),
             size=22, font=FONT_BLACK, color=ORANGE)
    add_bullets(slide, right_bullets, Inches(7.0), top + Inches(0.6),
                col_w, Inches(5), size=size)
